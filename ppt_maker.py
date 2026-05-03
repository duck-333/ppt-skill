#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT Maker (工业级克隆协议版)

核心约束（必须遵守）：
- 严禁新建 Shape：禁止 shapes.add_textbox / shapes.add_shape；除插入图片外，所有元素必须从模板页克隆出来。
- 精准定位：仅通过 Shape.name（选择窗格名字）定位并替换占位符；禁止坐标硬编码定位内容。
- 样式继承：替换文本必须保留模板的 font.name/font.size/font.color 等格式（通过 run 级别替换）。
- 分页逻辑：不在同一页叠加内容；每个单词=复制整张模板 Slide 2 -> 在新页上替换。
- 总起页：基于模板 Slide 1；按素材意项数 N 动态克隆/隐藏 SUB_ITEM_UNIT_X 组并更新序号与标题。
- 总结页：基于模板 Slide 3；按 N 网格克隆 SUMMARY_UNIT_TEMPLATE 组并填充。
- 发音按钮例外：详情页可追加“英/美”动作声音按钮，不修改模板，仅写入输出 PPT。
- AI 绘图预留：识别 PLACEHOLDER_AI_IMAGE，控制台打印该单词的 DALL·E 3 英文提示词。
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from pathlib import Path
from copy import deepcopy
from dataclasses import dataclass
from typing import Dict, Iterable, Iterator, List, Optional, Tuple
import re
import math
import json
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import asdict
from urllib.parse import quote
from urllib.request import Request, urlopen
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.media import Video
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Pt


TEMPLATE_PPTX = "master_template.pptx"
CONTENT_TXT = "content.txt"
OUTPUT_DEFAULT = "output_result.pptx"
AUDIO_CACHE_DIR = "audio_cache"
FFMPEG_CANDIDATES = (
    "ffmpeg",
    r"E:\downloads in E\SudaCaplayer\bin\Converter\x64\ffmpeg.exe",
    r"E:\downloads in E\SudaCaplayer\bin\Converter\ffmpeg.exe",
)

AUDIO_VARIANTS = (
    {
        "key": "uk",
        "label": "英",
        "culture": "en-GB",
        "youdao_type": "1",
        "language_lcids": {"809", "0809"},
        "accent_rgb": (21, 94, 239),
    },
    {
        "key": "us",
        "label": "美",
        "culture": "en-US",
        "youdao_type": "0",
        "language_lcids": {"409", "0409"},
        "accent_rgb": (255, 0, 0),
    },
)


def _safe_stdout_utf8() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass


def _safe_filename_piece(text: str) -> str:
    s = re.sub(r"[^A-Za-z0-9_-]+", "_", (text or "").strip().lower())
    return s.strip("_") or "word"


def _find_ffmpeg() -> Optional[str]:
    for candidate in FFMPEG_CANDIDATES:
        found = shutil.which(candidate)
        if found:
            return found
        p = Path(candidate)
        if p.exists():
            return str(p)
    return None


def _convert_mp3_to_action_wav(mp3_path: Path, wav_path: Path, warnings: set[str]) -> Optional[Path]:
    """
    PowerPoint action sounds are most reliable as PCM WAV. MP3 may embed, but often
    clicks without audible playback on some Office builds.
    """
    if wav_path.exists() and wav_path.stat().st_size > 1024:
        return wav_path
    ffmpeg = _find_ffmpeg()
    if not ffmpeg:
        warnings.add("未找到 ffmpeg，无法把在线 MP3 转成 PowerPoint 动作声音所需的 WAV。")
        return None
    wav_path.parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        ffmpeg,
        "-y",
        "-hide_banner",
        "-loglevel",
        "error",
        "-i",
        str(mp3_path),
        "-acodec",
        "pcm_s16le",
        "-ar",
        "22050",
        "-ac",
        "1",
        str(wav_path),
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception as exc:
        warnings.add(f"MP3 转 WAV 失败：{mp3_path.name} ({exc})")
        return None
    return wav_path if wav_path.exists() and wav_path.stat().st_size > 1024 else None


def _normalize_lcid(value: str) -> str:
    s = (value or "").strip().lower().replace("0x", "")
    return s.lstrip("0") or s


def _sapi_voices() -> List[Dict[str, str]]:
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore

        pythoncom.CoInitialize()
        speaker = win32com.client.Dispatch("SAPI.SpVoice")
        voices: List[Dict[str, str]] = []
        for token in speaker.GetVoices():
            try:
                voices.append(
                    {
                        "id": str(token.Id),
                        "description": str(token.GetDescription()),
                        "language": _normalize_lcid(str(token.GetAttribute("Language"))),
                    }
                )
            except Exception:
                continue
        return voices
    except Exception:
        return []


def _pick_sapi_voice(variant: Dict[str, object]) -> Optional[Dict[str, str]]:
    target_lcids = {_normalize_lcid(str(x)) for x in variant["language_lcids"]}  # type: ignore[index]
    for voice in _sapi_voices():
        if voice.get("language") in target_lcids:
            return voice
    return None


def _download_word_audio(
    *,
    word: str,
    variant: Dict[str, object],
    cache_dir: Path,
    warnings: set[str],
) -> Optional[Path]:
    cache_dir.mkdir(parents=True, exist_ok=True)
    audio_path = cache_dir / f"{_safe_filename_piece(word)}_{variant['key']}.mp3"
    if audio_path.exists() and audio_path.stat().st_size > 1024:
        return audio_path

    url = f"https://dict.youdao.com/dictvoice?type={variant['youdao_type']}&audio={quote(word)}"
    try:
        req = Request(
            url,
            headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            },
        )
        with urlopen(req, timeout=15) as resp:
            data = resp.read()
            content_type = (resp.headers.get("Content-Type") or "").lower()
        if len(data) <= 1024 or "audio" not in content_type:
            warnings.add(f"{variant['culture']} 在线发音返回内容异常，将尝试本地系统语音兜底。")
            return None
        audio_path.write_bytes(data)
        return audio_path
    except Exception:
        warnings.add(f"{variant['culture']} 在线发音下载失败，将尝试本地系统语音兜底。")
        return None


def _synthesize_word_audio(
    *,
    word: str,
    variant: Dict[str, object],
    cache_dir: Path,
    warnings: set[str],
) -> Optional[Path]:
    """
    Generate one pronunciation wav using local Windows SAPI voices.
    If the requested accent voice is not installed, skip it instead of faking the accent.
    """
    voice_info = _pick_sapi_voice(variant)
    culture = str(variant["culture"])
    if voice_info is None:
        warnings.add(f"未找到 {culture} 系统语音，已跳过 {variant['label']} 发音按钮。")
        return None

    cache_dir.mkdir(parents=True, exist_ok=True)
    audio_path = cache_dir / f"{_safe_filename_piece(word)}_{variant['key']}.wav"
    if audio_path.exists() and audio_path.stat().st_size > 1024:
        return audio_path

    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore

        pythoncom.CoInitialize()
        speaker = win32com.client.Dispatch("SAPI.SpVoice")
        for token in speaker.GetVoices():
            if str(token.Id) == voice_info["id"]:
                speaker.Voice = token
                break

        stream = win32com.client.Dispatch("SAPI.SpFileStream")
        # 3 = SSFMCreateForWrite. Keeping the numeric constant avoids requiring makepy.
        stream.Open(str(audio_path), 3, False)
        speaker.AudioOutputStream = stream
        speaker.Rate = -1
        speaker.Volume = 100
        speaker.Speak(word)
        stream.Close()
        return audio_path if audio_path.exists() and audio_path.stat().st_size > 1024 else None
    except Exception as exc:
        warnings.add(f"{culture} 发音音频生成失败：{word} ({exc})")
        return None


def _get_word_audio(
    *,
    word: str,
    variant: Dict[str, object],
    cache_dir: Path,
    warnings: set[str],
) -> Optional[Path]:
    mp3_path = _download_word_audio(word=word, variant=variant, cache_dir=cache_dir, warnings=warnings)
    if mp3_path is not None:
        wav_path = mp3_path.with_suffix(".wav")
        converted = _convert_mp3_to_action_wav(mp3_path, wav_path, warnings)
        if converted is not None:
            return converted
    return _synthesize_word_audio(word=word, variant=variant, cache_dir=cache_dir, warnings=warnings)


@dataclass
class WordItem:
    word: str
    phonetic: str = ""
    definition: str = ""
    analysis: str = ""


@dataclass
class SubMeaning:
    title: str
    words: List[WordItem]


@dataclass
class ParsedMaterial:
    top_category: str
    root_name: str
    root_logic: str
    meanings: List[SubMeaning]


_RE_WORD_LINE = re.compile(
    r"^(?P<word>[A-Za-z][A-Za-z\-]*)\s*(?:(?P<phonetic>/[^/]+/)\s*[，,]?\s*|[，,]\s*)(?P<rest>.*)$"
)
_RE_TOP_CATEGORY_PREFIX = re.compile(r"^\s*[（(]\s*[一二三四五六七八九十百千万0-9]+\s*[）)]\s*")
_RE_TOP_LESSON_PREFIX = re.compile(r"^\s*第\s*\d+\s*课\s*")
_RE_TOP_NUMBER_PREFIX = re.compile(r"^\s*\d+\s*[\.、．]\s*")
_RE_WORD_START = re.compile(r"^[A-Za-z][A-Za-z\-]*\s*/[^/]+/")
_POS_PATTERN = r"(?:n|v|adj|adv|prep|conj|pron|int|interj|num|art)\."
_RE_HAS_POS = re.compile(rf"\b{_POS_PATTERN}")
_CIRCLED_NUMS = set("①②③④⑤⑥⑦⑧⑨")


def _normalize_top_category(raw: str) -> str:
    s = (raw or "").strip().lstrip("\ufeff")
    if not s:
        return s
    s = _RE_TOP_CATEGORY_PREFIX.sub("", s).strip()
    s = _RE_TOP_LESSON_PREFIX.sub("", s).strip()
    s = _RE_TOP_NUMBER_PREFIX.sub("", s).strip()
    return s.rstrip("：:").strip()


def _compact_top_category_for_corner(normalized: str) -> str:
    """
    The small top-left corner label should identify the letter course only,
    e.g. "字母b--树干" -> "B", while root/body titles keep full meaning text.
    """
    s = (normalized or "").strip()
    m = re.match(r"^字母\s*([A-Za-z])", s)
    if m:
        return m.group(1).upper()
    return s


def _is_heading_line(line: str) -> bool:
    s = (line or "").strip()
    if not s:
        return False
    return bool(
        s.startswith("【释义】")
        or s.startswith("词根")
        or s[:1] in _CIRCLED_NUMS
        or re.match(r"^[（(]?[一二三四五六七八九十0-9]+[）)\.、．]", s)
        or re.match(r"^\d+\s*[\.\-、．]", s)
    )


def _should_merge_docx_line(prev: str, cur: str) -> bool:
    if not prev or not cur:
        return False
    if _is_heading_line(cur) or _RE_WORD_START.match(cur):
        return False
    if prev.startswith("【释义】"):
        return True
    if _RE_WORD_START.match(prev) and not _RE_HAS_POS.search(prev):
        return True
    if prev.endswith(("，", "、", "；", ":", "：", "+", "（", "(")):
        return True
    return False


def _candidate_word_line(line: str) -> str:
    """
    只接受真正以英文单词开头的词条行，避免把“字母b，...”或“【释义】词语bare...”
    这类说明句误拆成没有音标的单词页。
    """
    s = (line or "").strip().replace("\u00a0", " ")
    if not s or s.startswith("【释义】") or s.startswith(("字母", "词根")):
        return ""
    s = re.sub(r"^[\s•·\u2022\-—–*]+", "", s).strip()
    if re.match(r"^[A-Za-z][A-Za-z\-]*\s*(?:(?:/[^/]+/)\s*[，,]?|[，,])", s):
        return s
    return ""


def _clean_meaning_title(title: str) -> str:
    s = (title or "").strip().rstrip("：:")
    if s and s[0] in _CIRCLED_NUMS:
        s = s[1:].strip()
    s = re.sub(r"^\d+\s*[\.、．\-]\s*", "", s).strip()
    if "：" in s:
        before, after = s.split("：", 1)
        # Word 中偶尔会把标题和后续解释挤在同一段；标题只保留冒号前的概念名。
        if ("“" in before or before.startswith(("字母", "词根", "前缀"))) and len(after) > 0:
            s = before.strip()
    return s.rstrip("：:").strip()


def extract_content_from_docx(docx_path: Path) -> str:
    """
    从 docx 提取纯文本并做轻量清洗，输出可直接供 parse_content_to_tree 读取的文本。
    """
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    lines: List[str] = []

    with zipfile.ZipFile(str(docx_path), "r") as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)

    for para in root.findall(".//w:p", ns):
        texts = [(t.text or "") for t in para.findall(".//w:t", ns)]
        line = "".join(texts).strip()
        if not line:
            continue

        # 清理从文档中带出的噪声（图片旁边的单字母注记等）
        line = re.sub(r"^[aA]{1,3}(?=[\u4e00-\u9fff【（(])", "", line).strip()
        if not line:
            continue
        if re.fullmatch(r"[aA]+", line):
            continue
        if line in {"|", "丨", "习"}:
            continue

        if lines and _should_merge_docx_line(lines[-1], line):
            lines[-1] = f"{lines[-1]}{line}"
        else:
            lines.append(line)

    return ("\n".join(lines).strip() + "\n") if lines else ""


def parse_content_to_tree(text: str) -> ParsedMaterial:
    """
    解析为树状结构：大类 -> 多个子词根意项 -> 每个意项下的单词列表。

    面向当前 content.txt 的“工业可用”解析：
    - 大类：第一段标题行（如 （一）字母A——鹰-鹰）
    - ROOT_NAME：第一个形如 “（1）词根 ...” 行里取 “词根 av-” 或 “字母A——鹰-鹰”
    - ROOT_LOGIC：紧随其后的说明句（或首段描述）
    - 意项：形如 ①/②/③ 或 “①词根 av- ...” 的段落标题
    - 单词行：形如 “aviate /.../，即 ...，v. ...。”
    """
    lines = [ln.rstrip() for ln in (text or "").splitlines()]
    lines = [ln for ln in lines if ln.strip()]

    top_category = ""
    root_name = ""
    root_logic = ""

    # 1) top_category
    for ln in lines:
        if ln.startswith("（") and "）" in ln:
            top_category = ln.strip()
            break
    if not top_category and lines:
        top_category = lines[0].strip()
    full_top_category = _normalize_top_category(top_category)
    top_category = _compact_top_category_for_corner(full_top_category)

    # 2) root_name
    # 对“字母课”优先使用顶层标题，避免被后文某个词根误覆盖
    if full_top_category.startswith("字母"):
        root_name = full_top_category
    else:
        for ln in lines:
            if "词根" in ln and ("“" in ln or "：" in ln or ":" in ln):
                m = re.search(r"词根\s*([^\s，：:]+)", ln)
                if m:
                    root_name = f"词根 {m.group(1)}"
                    break
        if not root_name:
            # fallback: 重用顶层标题
            root_name = full_top_category.strip() if full_top_category else "ROOT"

    # 3) root_logic：取第一段解释句（在 root_name 段落附近）
    for i, ln in enumerate(lines):
        if ln.startswith("字母") or ("词根" in ln and "，" in ln):
            # 取后面 1~2 行作为逻辑摘要
            root_logic = ln.strip()
            if i + 1 < len(lines) and len(root_logic) < 60:
                nxt = lines[i + 1].strip()
                if not nxt.startswith("【释义】") and not nxt.startswith("①") and not nxt.startswith("②") and not nxt.startswith("③"):
                    root_logic = f"{root_logic}\n{nxt}"
            break
    if not root_logic:
        root_logic = " ".join(lines[:2]).strip()

    # 4) meanings + words
    meanings: List[SubMeaning] = []
    current: Optional[SubMeaning] = None

    def start_meaning(title: str) -> None:
        nonlocal current
        current = SubMeaning(title=title.strip(), words=[])
        meanings.append(current)

    for ln in lines:
        # 意项标题：①/②/③ 开头
        if (
            ln[:1] in _CIRCLED_NUMS
            or re.match(r"^[①②③④⑤⑥⑦⑧⑨]\\s*词根", ln)
            or re.match(r"^\d+\s*[\.、．\-]\s*字母", ln)
        ):
            start_meaning(_clean_meaning_title(ln))
            continue

        # 释义说明归并到上一条真实词条，避免生成“无音标重复页”。
        if ln.strip().startswith("【释义】"):
            if current is not None and current.words:
                last = current.words[-1]
                explanation = ln.strip()
                last.analysis = f"{last.analysis}\n{explanation}".strip() if last.analysis else explanation
            continue

        # 单词行
        ln_norm = _candidate_word_line(ln)
        m = _RE_WORD_LINE.match(ln_norm) if ln_norm else None
        if m:
            if current is None:
                start_meaning("意项")
            word = (m.group("word") or "").strip()
            phonetic = (m.group("phonetic") or "").strip()
            phonetic = phonetic.strip("/") if phonetic.startswith("/") else phonetic
            rest = (m.group("rest") or "").strip()
            if not phonetic and not _RE_HAS_POS.search(rest):
                # 例如资料说明中的 “baris，一次可装载...” 不是词条，避免生成空音标页。
                continue

            definition = ""
            analysis = ""
            # rest 常见：“即 ...，v. ...。”
            if "即" in rest:
                parts = rest.split("即", 1)
                after = parts[1].strip("，。 ")
                # 定义通常在最后一个 “，n.” / “，v.” 之后
                # 解析策略：保留全文为 analysis，definition 尝试抓最后一个 “，x.” 后面的部分
                analysis = f"即 {after}"
                def_m = re.search(rf"[，,]\s*({_POS_PATTERN}\s*[^。]+)", rest)
                if def_m:
                    definition = def_m.group(1).strip().rstrip("。")
                else:
                    # 兜底：词性标记可能不在逗号后（或 OCR 扰动）
                    def_m2 = re.search(rf"({_POS_PATTERN}\s*[^。]+)", rest)
                    if def_m2:
                        definition = def_m2.group(1).strip().rstrip("。")
            else:
                definition = rest.strip().rstrip("。")

            current.words.append(WordItem(word=word, phonetic=phonetic, definition=definition, analysis=analysis))

    # 清理空分组（例如“2.字母a-人头”这类章节标题，不直接承载单词）
    meanings = [m for m in meanings if m.words]

    # 兜底：若没识别到任何意项/单词，生成空结构，避免崩
    if not meanings:
        meanings = [SubMeaning(title="意项", words=[])]

    return ParsedMaterial(
        top_category=top_category,
        root_name=root_name,
        root_logic=root_logic,
        meanings=meanings,
    )


def _iter_shapes(shapes) -> Iterator[object]:
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(sh.shapes)


def find_shape_by_name(slide, name: str) -> Optional[object]:
    for sh in _iter_shapes(slide.shapes):
        if getattr(sh, "name", None) == name:
            return sh
    return None


def _replace_text_preserve_format(shape, mapping: Dict[str, str]) -> None:
    """
    run 级别替换，保留模板的字体/字号/颜色等。
    """
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        paragraph_changed = False
        for r in p.runs:
            if not r.text:
                continue
            t = r.text
            for k, v in mapping.items():
                if k in t:
                    t = t.replace(k, v)
            if t != r.text:
                r.text = t
                paragraph_changed = True

        # 兜底：处理占位符被拆分到多个 run 的情况
        full_text = "".join([r.text or "" for r in p.runs])
        if full_text:
            replaced = full_text
            for k, v in mapping.items():
                if k in replaced:
                    replaced = replaced.replace(k, v)
            if replaced != full_text and not paragraph_changed:
                if p.runs:
                    p.runs[0].text = replaced
                    for rr in p.runs[1:]:
                        rr.text = ""


def _replace_tokens_on_slide(slide, mapping: Dict[str, str]) -> None:
    for sh in _iter_shapes(slide.shapes):
        _replace_text_preserve_format(sh, mapping)


def _get_run_font_pt(shape, fallback: float = 20.0) -> float:
    if not hasattr(shape, "text_frame"):
        return fallback
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            try:
                if r.font and r.font.size:
                    return float(r.font.size.pt)
            except Exception:
                continue
    return fallback


def _estimate_text_height_emu(text: str, *, font_pt: float, box_width_emu: int) -> int:
    font_pt = max(12.0, float(font_pt))
    max_units_per_line = max(4.0, box_width_emu / max(1, font_pt * 12700))

    def char_units(ch: str) -> float:
        code = ord(ch)
        if ch.isspace():
            return 0.35
        if "\u4e00" <= ch <= "\u9fff":
            return 0.88
        if ch in "，。；：、“”‘’（）()[]{}【】+-/\\|":
            return 0.45
        if ch in ".,;:'\"`!iIljrtf":
            return 0.32
        if ch.isascii():
            return 0.55
        # IPA and other Latin-like symbols are usually narrower than CJK.
        if code < 0x0300 or 0x1D00 <= code <= 0x1DFF:
            return 0.55
        return 0.75

    logical_lines = 0
    for ln in (text or "").split("\n"):
        ln = ln.strip()
        if not ln:
            logical_lines += 1
            continue
        units = sum(char_units(ch) for ch in ln)
        logical_lines += max(1, math.ceil(units / max_units_per_line))

    line_h = int(font_pt * 1.35 * 12700)
    return max(line_h, logical_lines * line_h)


def _estimate_text_width_emu(text: str, *, font_pt: float) -> int:
    font_pt = max(12.0, float(font_pt))

    def char_units(ch: str) -> float:
        code = ord(ch)
        if ch.isspace():
            return 0.35
        if "\u4e00" <= ch <= "\u9fff":
            return 0.88
        if ch in "，。；：、“”‘’（）()[]{}【】+-/\\|":
            return 0.45
        if ch in ".,;:'\"`!iIljrtf":
            return 0.32
        if ch.isascii():
            return 0.55
        if code < 0x0300 or 0x1D00 <= code <= 0x1DFF:
            return 0.55
        return 0.75

    units = sum(char_units(ch) for ch in (text or ""))
    return int(units * font_pt * 12700)


def _set_text_wrap(shape, *, wrap: bool = True) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    try:
        tf.word_wrap = wrap
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        body_pr = shape.element.find(".//a:bodyPr", namespaces=shape.element.nsmap)
        if body_pr is not None:
            body_pr.set("wrap", "square" if wrap else "none")
    except Exception:
        pass


def _set_shape_font_size(shape, font_pt: float) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            try:
                r.font.size = Pt(float(font_pt))
            except Exception:
                continue


def _set_shape_text_color(shape, color: RGBColor) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            try:
                r.font.color.rgb = color
            except Exception:
                continue


def _shape_contains(outer, inner, *, tolerance_emu: int = 30000) -> bool:
    if outer is None or inner is None:
        return False
    ox1, oy1 = int(outer.left), int(outer.top)
    ox2, oy2 = ox1 + int(outer.width), oy1 + int(outer.height)
    ix1, iy1 = int(inner.left), int(inner.top)
    ix2, iy2 = ix1 + int(inner.width), iy1 + int(inner.height)
    return (
        ox1 - tolerance_emu <= ix1
        and oy1 - tolerance_emu <= iy1
        and ox2 + tolerance_emu >= ix2
        and oy2 + tolerance_emu >= iy2
    )


def _find_word_card_container(slide, inner_shapes: Iterable[object]) -> Optional[object]:
    named = find_shape_by_name(slide, "WORD_CARD_BG")
    if named is not None:
        return named
    inners = [sh for sh in inner_shapes if sh is not None]
    candidates: List[Tuple[int, object]] = []
    for sh in slide.shapes:
        if sh in inners:
            continue
        if not all(hasattr(sh, attr) for attr in ("left", "top", "width", "height")):
            continue
        if all(_shape_contains(sh, inner) for inner in inners):
            area = int(sh.width) * int(sh.height)
            candidates.append((area, sh))
    if not candidates:
        return None
    return sorted(candidates, key=lambda item: item[0])[0][1]


def _fit_text_box_inside_card(
    shape,
    text: str,
    *,
    card_shape=None,
    preferred_font_pt: Optional[float] = None,
    min_font_pt: float = 18.0,
    bottom_margin_emu: int = 250000,
) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return

    _set_text_wrap(shape, wrap=True)

    if card_shape is not None:
        right_limit = int(card_shape.left) + int(card_shape.width) - int(bottom_margin_emu)
        max_width = max(1000000, right_limit - int(shape.left))
        if int(shape.width) > max_width:
            shape.width = int(max_width)

        bottom_limit = int(card_shape.top) + int(card_shape.height) - int(bottom_margin_emu)
        max_height = max(int(shape.height), bottom_limit - int(shape.top))
        shape.height = int(max_height)

    font_pt = float(preferred_font_pt) if preferred_font_pt is not None else _get_run_font_pt(shape, fallback=20.0)
    fitted_pt = font_pt
    if preferred_font_pt is not None:
        candidates = list(range(int(round(font_pt)), int(round(min_font_pt)) - 1, -2))
        if candidates[-1] != int(round(min_font_pt)):
            candidates.append(int(round(min_font_pt)))
    else:
        candidates = [font_pt, 22, 20, 18, 16, 14]

    for candidate in candidates:
        candidate = min(float(candidate), fitted_pt)
        if candidate < min_font_pt:
            continue
        need_h = _estimate_text_height_emu(text, font_pt=candidate, box_width_emu=int(shape.width))
        if need_h <= int(shape.height):
            fitted_pt = candidate
            break
        fitted_pt = candidate

    _set_shape_font_size(shape, fitted_pt)


def _arrange_word_detail_text_blocks(
    *,
    word_shape,
    phonetic_box,
    definition_box,
    analysis_box,
    card_shape,
) -> None:
    """Place long analysis above the short Chinese definition inside the word card."""
    if definition_box is None or analysis_box is None:
        return
    if not all(hasattr(sh, attr) for sh in (definition_box, analysis_box) for attr in ("left", "top", "width", "height")):
        return

    gap = 180000
    side_margin = 250000
    bottom_margin = 300000

    left = int(analysis_box.left)
    if word_shape is not None and hasattr(word_shape, "left"):
        left = int(word_shape.left)

    if card_shape is not None and all(hasattr(card_shape, attr) for attr in ("left", "top", "width", "height")):
        card_right = int(card_shape.left) + int(card_shape.width) - side_margin
        card_bottom = int(card_shape.top) + int(card_shape.height) - bottom_margin
    else:
        card_right = max(int(analysis_box.left) + int(analysis_box.width), int(definition_box.left) + int(definition_box.width))
        card_bottom = int(analysis_box.top) + int(analysis_box.height)

    usable_width = max(int(analysis_box.width), card_right - left)
    analysis_box.left = left
    definition_box.left = left
    analysis_box.width = usable_width
    definition_box.width = usable_width

    if phonetic_box is not None and hasattr(phonetic_box, "top") and hasattr(phonetic_box, "height"):
        analysis_top = int(phonetic_box.top) + int(phonetic_box.height) + 260000
    else:
        analysis_top = int(definition_box.top)

    definition_font = 36.0
    definition_needed = _estimate_text_height_emu(
        definition_box.text if hasattr(definition_box, "text") else "",
        font_pt=definition_font,
        box_width_emu=int(definition_box.width),
    )
    definition_height = max(int(definition_box.height), definition_needed + 120000)

    available_analysis_height = card_bottom - analysis_top - gap - definition_height
    if available_analysis_height < 900000:
        # Keep the definition readable even on crowded cards, but still reserve a usable
        # upper region for the explanatory paragraph.
        available_analysis_height = max(650000, card_bottom - analysis_top - gap - int(definition_box.height))
        definition_height = max(int(definition_box.height), card_bottom - analysis_top - gap - available_analysis_height)

    analysis_box.top = analysis_top
    analysis_box.height = max(650000, available_analysis_height)
    definition_box.top = int(analysis_box.top) + int(analysis_box.height) + gap
    definition_box.height = max(int(definition_box.height), definition_height)


def _fit_word_name(shape) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    _set_text_wrap(shape, wrap=False)
    try:
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _new_sound_action_button_xml(
    *,
    shape_id: int,
    shape_name: str,
    audio_rid: str,
    audio_name: str,
    label: str,
    accent_rgb: Tuple[int, int, int],
    x: int,
    y: int,
    cx: int,
    cy: int,
):
    accent_hex = "%02X%02X%02X" % tuple(int(v) for v in accent_rgb)
    return parse_xml(
        (
            "<p:sp %s>\n"
            "  <p:nvSpPr>\n"
            '    <p:cNvPr id="%d" name="%s">\n'
            '      <a:hlinkClick action="ppaction://noaction">\n'
            '        <a:snd r:embed="%s" name="%s"/>\n'
            "      </a:hlinkClick>\n"
            "    </p:cNvPr>\n"
            "    <p:cNvSpPr/>\n"
            "    <p:nvPr/>\n"
            "  </p:nvSpPr>\n"
            "  <p:spPr>\n"
            "    <a:xfrm>\n"
            '      <a:off x="%d" y="%d"/>\n'
            '      <a:ext cx="%d" cy="%d"/>\n'
            "    </a:xfrm>\n"
            '    <a:prstGeom prst="roundRect">\n'
            "      <a:avLst/>\n"
            "    </a:prstGeom>\n"
            "    <a:solidFill>\n"
            '      <a:srgbClr val="%s"/>\n'
            "    </a:solidFill>\n"
            '    <a:ln w="25400">\n'
            "      <a:solidFill>\n"
            '        <a:srgbClr val="FFFFFF"/>\n'
            "      </a:solidFill>\n"
            "    </a:ln>\n"
            "  </p:spPr>\n"
            "  <p:style>\n"
            "    <a:lnRef idx=\"2\"><a:schemeClr val=\"accent1\"/></a:lnRef>\n"
            "    <a:fillRef idx=\"1\"><a:schemeClr val=\"accent1\"/></a:fillRef>\n"
            "    <a:effectRef idx=\"0\"><a:schemeClr val=\"accent1\"/></a:effectRef>\n"
            "    <a:fontRef idx=\"minor\"><a:schemeClr val=\"lt1\"/></a:fontRef>\n"
            "  </p:style>\n"
            "  <p:txBody>\n"
            '    <a:bodyPr wrap="none" anchor="ctr">\n'
            "      <a:spAutoFit/>\n"
            "    </a:bodyPr>\n"
            "    <a:lstStyle/>\n"
            '    <a:p>\n'
            '      <a:pPr algn="ctr"/>\n'
            '      <a:r>\n'
            '        <a:rPr lang="zh-CN" sz="3400" b="1">\n'
            '          <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>\n'
            '          <a:latin typeface="Microsoft YaHei"/>\n'
            '          <a:ea typeface="Microsoft YaHei"/>\n'
            "        </a:rPr>\n"
            "        <a:t>%s</a:t>\n"
            "      </a:r>\n"
            "    </a:p>\n"
            "  </p:txBody>\n"
            "</p:sp>\n"
        )
        % (
            nsdecls("a", "p", "r"),
            int(shape_id),
            escape(shape_name),
            audio_rid,
            escape(audio_name),
            int(x),
            int(y),
            int(cx),
            int(cy),
            accent_hex,
            escape(label),
        )
    )


def _add_embedded_audio_button(
    *,
    slide,
    audio_path: Path,
    label: str,
    accent_rgb: Tuple[int, int, int],
    x: int,
    y: int,
    cx: int,
    cy: int,
    shape_name: str,
) -> None:
    mime_type = "audio/mpeg" if audio_path.suffix.lower() == ".mp3" else "audio/wav"
    media = Video.from_path_or_file_like(str(audio_path), mime_type)
    media_part = slide.part._package.get_or_add_media_part(media)
    audio_rid = slide.part.relate_to(media_part, RT.AUDIO)

    shape_id = int(slide.shapes._next_shape_id)
    button = _new_sound_action_button_xml(
        shape_id=shape_id,
        shape_name=shape_name,
        audio_rid=audio_rid,
        audio_name=audio_path.name,
        label=label,
        accent_rgb=accent_rgb,
        x=int(x),
        y=int(y),
        cx=int(cx),
        cy=int(cy),
    )
    slide.shapes._spTree.append(button)


def _add_pronunciation_buttons(
    *,
    slide,
    word: str,
    word_shape,
    card_shape,
    audio_cache_dir: Path,
    warnings: set[str],
) -> None:
    if word_shape is None:
        return

    available: List[Tuple[Dict[str, object], Path]] = []
    for variant in AUDIO_VARIANTS:
        audio_path = _get_word_audio(
            word=word,
            variant=variant,
            cache_dir=audio_cache_dir,
            warnings=warnings,
        )
        if audio_path is not None:
            available.append((variant, audio_path))
    if not available:
        return

    btn_w = 650000
    btn_h = 650000
    gap = 260000
    total_w = len(available) * btn_w + max(0, len(available) - 1) * gap
    margin = 330000

    if card_shape is not None:
        right_limit = int(card_shape.left) + int(card_shape.width) - margin
        top_limit = int(card_shape.top) + 180000
    else:
        right_limit = int(word_shape.left) + int(word_shape.width)
        top_limit = max(0, int(word_shape.top) - btn_h - gap)

    word_font = _get_run_font_pt(word_shape, fallback=68.0)
    word_end = int(word_shape.left) + _estimate_text_width_emu(word, font_pt=word_font)
    x = max(int(word_shape.left), right_limit - total_w)
    y = int(word_shape.top) + 120000

    # If a long word reaches the button area, lift the buttons into the card's upper margin.
    if word_end + 140000 > x:
        y = max(top_limit, int(word_shape.top) - btn_h - gap)

    for idx, (variant, audio_path) in enumerate(available):
        _add_embedded_audio_button(
            slide=slide,
            audio_path=audio_path,
            label=str(variant["label"]),
            accent_rgb=tuple(variant["accent_rgb"]),  # type: ignore[arg-type]
            x=x + idx * (btn_w + gap),
            y=y,
            cx=btn_w,
            cy=btn_h,
            shape_name=f"PRON_{str(variant['key']).upper()}_{word}",
        )


def _set_paragraph_font_size(shape, paragraph_index: int, font_pt: float) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    paragraphs = shape.text_frame.paragraphs
    if paragraph_index >= len(paragraphs):
        return
    for r in paragraphs[paragraph_index].runs:
        try:
            r.font.size = Pt(float(font_pt))
        except Exception:
            continue


def _text_frame_inner_width(shape) -> int:
    if shape is None or not hasattr(shape, "text_frame"):
        return int(shape.width) if shape is not None else 0
    tf = shape.text_frame
    try:
        left = int(tf.margin_left or 0)
        right = int(tf.margin_right or 0)
    except Exception:
        left = right = 0
    return max(300000, int(shape.width) - left - right)


def _summary_card_required_height(
    *,
    title: str,
    word_list: str,
    title_font_pt: float,
    words_font_pt: float,
    width_emu: int,
    vertical_padding_emu: int = 420000,
    paragraph_gap_emu: int = 160000,
) -> int:
    title_h = _estimate_text_height_emu(title, font_pt=title_font_pt, box_width_emu=width_emu)
    words_h = _estimate_text_height_emu(word_list, font_pt=words_font_pt, box_width_emu=width_emu)
    return int(title_h + words_h + vertical_padding_emu + paragraph_gap_emu)


def _fit_summary_card(
    card,
    *,
    title: str,
    word_list: str,
    slide_height_emu: int,
    min_height_emu: int,
) -> None:
    if card is None or not hasattr(card, "text_frame"):
        return

    _set_text_wrap(card, wrap=True)
    inner_w = _text_frame_inner_width(card)
    title_font = 28.0 if len(title) <= 16 else 24.0
    max_bottom = int(slide_height_emu) - 520000
    available_h = max(int(min_height_emu), max_bottom - int(card.top))

    chosen_words_font = 18.0
    chosen_h = available_h
    for words_font in [28.0, 26.0, 24.0, 22.0, 20.0, 18.0]:
        needed_h = _summary_card_required_height(
            title=title,
            word_list=word_list,
            title_font_pt=title_font,
            words_font_pt=words_font,
            width_emu=inner_w,
        )
        if needed_h <= available_h:
            chosen_words_font = words_font
            chosen_h = max(int(min_height_emu), int(needed_h))
            break

    card.height = int(chosen_h)
    _set_paragraph_font_size(card, 0, title_font)
    _set_paragraph_font_size(card, 1, chosen_words_font)


def _shape_set_name(shape_element, new_name: str) -> None:
    """
    通过 XML 修改 cNvPr/@name，让克隆后的 shape 仍可被 name 精准定位。
    """
    cNvPr = shape_element.find(".//p:cNvPr", namespaces=shape_element.nsmap)
    if cNvPr is not None:
        cNvPr.set("name", new_name)


def _set_group_transform_y(grp_element, new_y_emu: int) -> None:
    off = grp_element.find(".//a:xfrm/a:off", namespaces=grp_element.nsmap)
    if off is None:
        return
    off.set("y", str(int(new_y_emu)))


def _clone_group_on_slide(slide, group_shape, *, new_name: str, dy_emu: int) -> object:
    """
    克隆一个 Group 并插入到 slide，整体下移 dy_emu，并修改 name。
    注意：这是“克隆原件并替换”的唯一合法方式（不新建任何 Shape）。
    """
    grp_el = deepcopy(group_shape.element)
    # 下移：改 group xfrm 的 off.y（整体移动）
    try:
        off = grp_el.find(".//a:xfrm/a:off", namespaces=grp_el.nsmap)
        if off is not None and off.get("y") is not None:
            off.set("y", str(int(off.get("y")) + int(dy_emu)))
    except Exception:
        pass

    _shape_set_name(grp_el, new_name)
    slide.shapes._spTree.insert_element_before(grp_el, "p:extLst")
    # 返回新插入的 shape 对象：只能通过重新查找 name 获取
    return find_shape_by_name(slide, new_name)


def find_shape_in_group_by_name(group_shape, name: str) -> Optional[object]:
    """
    仅在指定 group 内部按 name 查找（局部定位，避免“全页替换”误命中）。
    """
    if group_shape is None or getattr(group_shape, "shape_type", None) != MSO_SHAPE_TYPE.GROUP:
        return None
    for sh in _iter_shapes(group_shape.shapes):
        if getattr(sh, "name", None) == name:
            return sh
    return None


def _clone_group_element_to_slide(*, slide, group_element, new_name: str) -> object:
    """
    从“原始 group XML”克隆到 slide，并命名为 new_name。
    """
    grp_el = deepcopy(group_element)
    _shape_set_name(grp_el, new_name)
    slide.shapes._spTree.insert_element_before(grp_el, "p:extLst")
    return find_shape_by_name(slide, new_name)


def _copy_slide_with_rels(*, src_prs: Presentation, dst_prs: Presentation, slide_index: int):
    """
    复制整张幻灯片（保留所有 Shape），并复制必要的关系（尤其图片 rId）。
    约束：不新建 Shape，仅新增一张 slide 并拷贝其 spTree。
    """
    if slide_index < 0 or slide_index >= len(src_prs.slides):
        raise IndexError(f"slide_index out of range: {slide_index}, total={len(src_prs.slides)}")

    src = src_prs.slides[slide_index]
    new_slide = dst_prs.slides.add_slide(src.slide_layout)

    # 清空 new_slide 上由 layout 自动产生的占位符（避免叠加）
    for shp in list(new_slide.shapes):
        try:
            new_slide.shapes._spTree.remove(shp.element)
        except Exception:
            pass

    # 复制关系：图片等需要重新 relate_to，得到新 rId，并更新 blip embed
    rid_map: Dict[str, str] = {}
    for rel in src.part.rels.values():
        if rel.is_external:
            continue
        # 跳过 layout/master 关系
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/slideMaster") or rel.reltype.endswith("/notesSlide"):
            continue
        # python-pptx：relate_to(target_part, reltype) -> new_rId
        try:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        except Exception:
            continue

    for shp in src.shapes:
        el = deepcopy(shp.element)
        # 更新图片 embed rid
        for blip in el.findall(".//a:blip", namespaces=el.nsmap):
            embed = blip.get(qn("r:embed"))
            if embed and embed in rid_map:
                blip.set(qn("r:embed"), rid_map[embed])
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    return new_slide


def build_high_fidelity_ppt(
    *,
    template_path: Path,
    content_path: Path,
    output_path: Path,
    enable_audio: bool = True,
    audio_cache_dir: Optional[Path] = None,
) -> Path:
    material = parse_content_to_tree(content_path.read_text(encoding="utf-8"))
    # 调试输出：打印解析后的树状 JSON，便于核对每个单词是否被解析为独立 Detail 节点
    print("[parse_content_to_tree] JSON:")
    print(json.dumps(asdict(material), ensure_ascii=False, indent=2))

    meanings = material.meanings
    n_meanings = len(meanings)
    audio_cache_dir = audio_cache_dir or (content_path.parent / AUDIO_CACHE_DIR)
    audio_warnings: set[str] = set()

    # 源模板（只读）
    src = Presentation(str(template_path))
    assert len(src.slides) >= 3, "错误：master_template.pptx 至少需要 3 页（总起、详情、总结）"

    # 输出文档：创建空白 presentation，并继承页面尺寸
    dst = Presentation()
    dst.slide_width = src.slide_width
    dst.slide_height = src.slide_height
    # 移除默认空白页
    try:
        sldIdLst = dst.slides._sldIdLst
        if len(sldIdLst) > 0:
            rid = sldIdLst[0].rId
            dst.part.drop_rel(rid)
            el = sldIdLst[0]
            el.getparent().remove(el)
    except Exception:
        pass

    # 1) 总起页：克隆模板 Slide 1 (index 0)
    slide_cover = _copy_slide_with_rels(src_prs=src, dst_prs=dst, slide_index=0)
    _replace_tokens_on_slide(
        slide_cover,
        {
            "{{ROOT_NAME}}": material.root_name,
            "{{ROOT_LOGIC}}": material.root_logic,
            "{{TOP_CATEGORY}}": material.top_category,
        },
    )

    # SUB_ITEM_UNIT 动态增减：模板有 1~3，超过 3 就克隆 SUB_ITEM_UNIT_3 继续向下排
    unit1 = find_shape_by_name(slide_cover, "SUB_ITEM_UNIT_1")
    unit2 = find_shape_by_name(slide_cover, "SUB_ITEM_UNIT_2")
    unit3 = find_shape_by_name(slide_cover, "SUB_ITEM_UNIT_3")
    units: List[object] = [u for u in [unit1, unit2, unit3] if u is not None]

    gap = 0
    unit_h = 0
    if unit1 is not None and unit2 is not None:
        unit_h = int(unit1.height)
        gap = int(unit2.top) - int(unit1.top) - int(unit1.height)
    elif unit3 is not None:
        unit_h = int(unit3.height)
        gap = int(unit_h * 0.25)

    # 克隆补足
    for idx in range(4, n_meanings + 1):
        if unit3 is None:
            break
        dy = (idx - 3) * int(unit_h + gap)
        _clone_group_on_slide(slide_cover, unit3, new_name=f"SUB_ITEM_UNIT_{idx}", dy_emu=dy)

    # 填充/隐藏（仅隐藏：把组的所有子 shape 文本清空）
    for i in range(1, max(3, n_meanings) + 1):
        grp = find_shape_by_name(slide_cover, f"SUB_ITEM_UNIT_{i}")
        if grp is None:
            continue
        if i <= n_meanings:
            title = meanings[i - 1].title
            # 子标题（兼容：SUB_ITEM_UNIT_* 既可能是 Group，也可能是单个 TextBox）
            _replace_text_preserve_format(find_shape_by_name(slide_cover, f"Text_Content_{min(i,3)}") or grp, {f"{{{{SUB_TITLE_{min(i,3)}}}}}": title})
            if getattr(grp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                # 若克隆后内部 name 未变，通过遍历 group 子形状按占位符替换
                for gsh in _iter_shapes(grp.shapes):
                    _replace_text_preserve_format(gsh, {
                        "{{SUB_TITLE_1}}": title,
                        "{{SUB_TITLE_2}}": title,
                        "{{SUB_TITLE_3}}": title,
                    })
                    # 序号：把“��”替换为 i
                    _replace_text_preserve_format(gsh, {"��": str(i)})
            else:
                _replace_text_preserve_format(grp, {
                    "{{SUB_TITLE_1}}": title,
                    "{{SUB_TITLE_2}}": title,
                    "{{SUB_TITLE_3}}": title,
                    "��": str(i),
                })
        else:
            # 多余模板 unit：清空占位符文本（视觉等价于隐藏）
            if getattr(grp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                for gsh in _iter_shapes(grp.shapes):
                    _replace_text_preserve_format(gsh, {
                        "{{SUB_TITLE_1}}": "",
                        "{{SUB_TITLE_2}}": "",
                        "{{SUB_TITLE_3}}": "",
                        "��": "",
                    })
            else:
                _replace_text_preserve_format(grp, {
                    "{{SUB_TITLE_1}}": "",
                    "{{SUB_TITLE_2}}": "",
                    "{{SUB_TITLE_3}}": "",
                    "��": "",
                })

    # 2) 详情页：每个单词复制整张模板 Slide 2
    for meaning in meanings:
        for w in meaning.words:
            slide_word = _copy_slide_with_rels(src_prs=src, dst_prs=dst, slide_index=1)
            phonetic_text = f"/{w.phonetic}/" if w.phonetic else ""
            _replace_tokens_on_slide(
                slide_word,
                {
                    "{{TOP_CATEGORY}}": material.top_category,
                    "{{SUB_TITLE}}": meaning.title,
                    "{{WORD}}": w.word,
                    "/{{PHONETIC}}/": phonetic_text,
                    "{{PHONETIC}}": w.phonetic,
                    "{{DEFINITION}}": w.definition,
                    "{{ROOT_ANALYSIS}}": w.analysis,
                },
            )

            word_name = find_shape_by_name(slide_word, "WORD_NAME")
            phonetic_box = find_shape_by_name(slide_word, "WORD_PHONETIC")
            definition_box = find_shape_by_name(slide_word, "WORD_DEFINITION")
            analysis_box = find_shape_by_name(slide_word, "WORD_ANALYSIS")

            # 详情页版式微调：主词变红，白卡内长解释自动换行并收在卡片范围内。
            _set_shape_text_color(word_name, RGBColor(255, 0, 0))
            _fit_word_name(word_name)
            _set_text_wrap(phonetic_box, wrap=False)

            card = _find_word_card_container(slide_word, [word_name, phonetic_box, definition_box, analysis_box])
            _arrange_word_detail_text_blocks(
                word_shape=word_name,
                phonetic_box=phonetic_box,
                definition_box=definition_box,
                analysis_box=analysis_box,
                card_shape=card,
            )
            _fit_text_box_inside_card(analysis_box, w.analysis, min_font_pt=16.0)
            _fit_text_box_inside_card(definition_box, w.definition, preferred_font_pt=36.0, min_font_pt=30.0)
            _set_shape_text_color(definition_box, RGBColor(255, 0, 0))
            if enable_audio:
                _add_pronunciation_buttons(
                    slide=slide_word,
                    word=w.word,
                    word_shape=word_name,
                    card_shape=card,
                    audio_cache_dir=audio_cache_dir,
                    warnings=audio_warnings,
                )

            # AI 图片提示词（仅打印，不创建 shape）
            if find_shape_by_name(slide_word, "PLACEHOLDER_AI_IMAGE") is not None:
                prompt = (
                    f"A clean, modern educational illustration of '{w.word}', "
                    f"flat vector style, minimal background, suitable for a language learning slide. "
                    f"No text, no watermark."
                )
                print(f"[DALL·E 3 prompt] {w.word}: {prompt}")

    # 3) 总结页：克隆模板 Slide 3 (index 2)。每页最多 3 个意项，给长单词列表留足高度。
    summary_page_size = 3
    summary_chunks = [meanings[i : i + summary_page_size] for i in range(0, n_meanings, summary_page_size)] or [[]]
    for chunk_index, meaning_chunk in enumerate(summary_chunks, start=1):
        slide_sum = _copy_slide_with_rels(src_prs=src, dst_prs=dst, slide_index=2)
        unit_tpl = find_shape_by_name(slide_sum, "SUMMARY_UNIT_TEMPLATE")
        if unit_tpl is None:
            continue

        # 保存“原始模板卡片”的 XML：后续每次都从这份原件克隆，确保占位符仍存在（数据解耦）
        unit_tpl_el_pristine = deepcopy(unit_tpl.element)

        base_left = int(unit_tpl.left)
        base_top = int(unit_tpl.top)
        unit_w = int(unit_tpl.width)
        unit_h = int(unit_tpl.height)
        gap_x = int(unit_w * 0.08)
        gap_y = int(unit_h * 0.18)

        try:
            slide_sum.shapes._spTree.remove(unit_tpl.element)
        except Exception:
            pass

        chunk_n = len(meaning_chunk)
        cols = max(1, min(3, chunk_n))
        available_w = int(dst.slide_width) - 2 * base_left
        fit_unit_w = int((available_w - (cols - 1) * gap_x) / cols)
        if chunk_n == 1:
            fit_unit_w = min(fit_unit_w, int(unit_w * 1.55))

        current_row = -1
        current_row_top = base_top
        current_row_h = unit_h

        for local_idx, meaning in enumerate(meaning_chunk, start=1):
            r = (local_idx - 1) // cols
            c = (local_idx - 1) % cols
            if r != current_row:
                if current_row >= 0:
                    current_row_top = int(current_row_top + current_row_h + gap_y)
                current_row = r
                current_row_h = unit_h

            target_left = base_left + c * (fit_unit_w + gap_x)
            target_top = current_row_top
            global_idx = (chunk_index - 1) * summary_page_size + local_idx

            new_name = f"SUMMARY_UNIT_{global_idx}"
            if getattr(unit_tpl, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                card = _clone_group_element_to_slide(
                    slide=slide_sum,
                    group_element=unit_tpl_el_pristine,
                    new_name=new_name,
                )
            else:
                el = deepcopy(unit_tpl_el_pristine)
                _shape_set_name(el, new_name)
                slide_sum.shapes._spTree.insert_element_before(el, "p:extLst")
                card = find_shape_by_name(slide_sum, new_name)

            if card is None:
                continue

            card.left = int(target_left)
            card.top = int(target_top)
            card.width = int(fit_unit_w)

            word_list = ", ".join([wi.word for wi in meaning.words])

            if getattr(card, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                subtitle_shape = find_shape_in_group_by_name(card, "SUMMARY_SUBTITLE")
                wordlist_shape = find_shape_in_group_by_name(card, "SUMMARY_WORD_LIST")
                if subtitle_shape is not None:
                    _replace_text_preserve_format(subtitle_shape, {"{{SUB_TITLE}}": meaning.title})
                if wordlist_shape is not None:
                    _replace_text_preserve_format(wordlist_shape, {"{{WORD_LIST}}": word_list})
            else:
                _replace_text_preserve_format(card, {"{{SUB_TITLE}}": meaning.title, "{{WORD_LIST}}": word_list})
                _fit_summary_card(
                    card,
                    title=meaning.title,
                    word_list=word_list,
                    slide_height_emu=int(dst.slide_height),
                    min_height_emu=unit_h,
                )
                current_row_h = max(current_row_h, int(card.height))

    for warning in sorted(audio_warnings):
        print(f"[pronunciation] {warning}")

    # 保存（避免被 PowerPoint 锁）
    candidates: List[Path] = [output_path]
    candidates.append(output_path.with_name(f"{output_path.stem}_alt{output_path.suffix}"))
    # 继续尝试更多不冲突的文件名（即使 PPT 正在打开也能保存成功）
    for i in range(1, 20):
        candidates.append(output_path.with_name(f"{output_path.stem}_{i:02d}{output_path.suffix}"))

    last_err: Optional[Exception] = None
    for p in candidates:
        try:
            dst.save(str(p))
            return p
        except PermissionError as e:
            last_err = e
            continue
    raise last_err  # type: ignore[misc]


if __name__ == "__main__":
    _safe_stdout_utf8()
    base = Path(__file__).resolve().parent

    parser = argparse.ArgumentParser(description="Generate PPT from content.txt or a DOCX source.")
    parser.add_argument("--template", default=str(base / TEMPLATE_PPTX), help="Path to master template pptx")
    parser.add_argument("--content", default=str(base / CONTENT_TXT), help="Path to content.txt")
    parser.add_argument("--docx", default=None, help="Optional path to source .docx; when set, auto-extract to content")
    parser.add_argument("--output", default=None, help="Optional output .pptx path")
    parser.add_argument("--no-audio", action="store_true", help="Disable pronunciation audio buttons")
    parser.add_argument(
        "--audio-cache",
        default=None,
        help="Optional directory for generated pronunciation wav cache",
    )
    args = parser.parse_args()

    tpl = Path(args.template)
    content = Path(args.content)

    if args.docx:
        docx_path = Path(args.docx)
        extracted = extract_content_from_docx(docx_path)
        content.write_text(extracted, encoding="utf-8")
        print(f"已从 DOCX 提取内容：{docx_path}")
        print(f"已写入内容文件：{content}")

    if args.output:
        output_path = Path(args.output)
    elif args.docx:
        output_path = base / f"{Path(args.docx).stem}.pptx"
    else:
        output_path = base / OUTPUT_DEFAULT

    out = build_high_fidelity_ppt(
        template_path=tpl,
        content_path=content,
        output_path=output_path,
        enable_audio=not args.no_audio,
        audio_cache_dir=Path(args.audio_cache) if args.audio_cache else None,
    )
    print(f"已生成：{out.resolve()}")

