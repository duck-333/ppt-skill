# -*- coding: utf-8 -*-
"""
从「a-鹰 第一版.pptx」按页克隆版式，生成「词根 au- 鸟」主题 PPT。
遵循 style_guide：元件来自原稿、仅替换文字；不新建简陋白底页。
"""
from __future__ import annotations

import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

TEMPLATE_NAME = "a-鹰 第一版.pptx"
OUTPUT_NAME = "output_au_bird.pptx"

# 输出幻灯片顺序 = 依次从模板复制的源页索引（对应先前与用户约定的拆页）
SOURCE_SLIDE_ORDER = [0, 1, 8, 2, 4, 5, 6, 7, 3]

# 每一页：{ 形状在模板中的枚举下标: 新文本 }（下标与模板 dump 一致）
REPLACEMENTS: list[dict[int, str]] = [
    # 0 起源/记忆点（模板第 0 页）
    {
        4: "词根 au",
        5: "词根 au - 鸟",
        6: "观鸟以占：古罗马人从飞鸟轨迹读取吉兆。",
        8: "au- 鸟",
        11: "1. au-鸟（占卜）",
        12: "2. austr-南风 · auto-自主",
    },
    # 1 释义/教学（模板第 1 页）
    {
        4: "词根 au",
        5: "词根 au - 鸟",
        6: "aul- 通 av-「鸟」：构词记忆锚点。",
        8: "词根 aul- 通 av-",
        9: "【释义】aul- 通词根 av-「鸟」，字母 ul 通 V。",
        11: "词根 au- 家族",
        13: "① au-「鸟」",
        14: "② austr-「南风；南方」",
        15: "③ auto-「自己；自主」",
    },
    # 2 三列总览（模板第 8 页）
    {
        3: "词根 au",
        7: "词根 au- 鸟",
        6: "augur\ninaugurate\ninauguration\nauspice",
        9: "词根 austr- 南风；南方",
        8: "Auster\naustral\nAustralia",
        12: "词根 auto- 自己；自主",
        11: "automobile\nauto\nautonomy",
    },
    # 3 augur + inaugurate（模板第 2 页）
    {
        3: "词根 au",
        12: "词根 au- 鸟",
        7: "augur /ˈɔːɡə(r)/",
        8: "n.占卜师；v.占卜；预言。",
        6: "即 au「鸟」+ gu通w「看」+ (o)r表人（观鸟占卜者）",
        9: "inaugurate /ɪˈnɔːɡjəreɪt/",
        10: "v.举行就职典礼；开创。",
        11: "【释义】古罗马重大仪式前先占卜；in-使 + augur + -ate",
    },
    # 4 Inauguration + auspice（模板第 4 页）
    {
        3: "词根 au",
        12: "词根 au- 鸟",
        7: "Inauguration /ɪˌnɔːɡjəˈreɪʃn/",
        8: "n.就职典礼；开幕式；开创。",
        6: "即 inaugurate「举行就职典礼」+ -ion 名词后缀",
        9: "auspice /ˈɔːspɪs/",
        10: "n.吉兆；赞助；主办。",
        11: "即 au「鸟」+ spic「看」+ e（观鸟占兆）",
    },
    # 5 austr-（模板第 5 页）
    {
        3: "词根 au",
        12: "词根 austr- 南风；南方",
        7: "Auster /ˈɔːstə(r)/",
        8: "n.南风神（奥斯忒耳）。",
        6: "【释义】罗马神话四大风神之一；翅翼男人形象；派生 austr-",
        9: "austral /ˈɒstrəl/",
        10: "adj.南方的；南风的。",
        11: "即 austr「南方」+ -al；Australia「南方大陆」",
    },
    # 6 auto- 引入 + automobile（模板第 6 页）
    {
        3: "词根 au",
        8: "词根 auto-：自己；自主\n（如鸟自飞，从心所欲）",
        5: "automobile /ˈɔːtəməbiːl/",
        6: "n.汽车。",
        7: "即 auto「自己」+ mobile「移动」",
    },
    # 7 auto + autonomy（模板第 7 页）
    {
        3: "词根 au",
        12: "词根 auto- 自己；自主",
        7: "auto /ˈɔːtəʊ/",
        8: "n.汽车（automobile 缩写）。",
        6: "即 automobile 的缩写。",
        9: "autonomy /ɔːˈtɒnəmi/",
        10: "n.自治；自主。",
        11: "即 auto「自己」+ nom「管理」+ -y",
    },
    # 8 三词柱复盘（模板第 3 页）
    {
        3: "词根 au",
        20: "词根 au- 鸟 · 复盘",
        6: "augur",
        7: "n./v.占卜；预言",
        8: "/ˈɔːɡə(r)/",
        5: "au「鸟」+ 看 + 人",
        11: "Auster",
        12: "南风神",
        13: "/ˈɔːstə(r)/",
        10: "翅翼风神 · austral",
        16: "autonomy",
        17: "自治；自主",
        18: "/ɔːˈtɒnəmi/",
        15: "auto + nom「管理」",
        19: "如鸟自飞 · 从心所欲",
    },
]


def _strip_default_slide(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    if len(sld_id_lst) == 0:
        return
    r_id = sld_id_lst[0].rId
    prs.part.drop_rel(r_id)
    el = sld_id_lst[0]
    el.getparent().remove(el)


def duplicate_slide_into(src: Presentation, src_index: int, dst: Presentation):
    blank = dst.slide_layouts[6]
    new_slide = dst.slides.add_slide(blank)
    source_slide = src.slides[src_index]
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return new_slide


def apply_replacements(slide, mapping: dict[int, str]) -> None:
    for idx, sh in enumerate(slide.shapes):
        if idx not in mapping:
            continue
        if not hasattr(sh, "text_frame"):
            continue
        sh.text = mapping[idx]


def build(output_path: Path | None = None) -> Path:
    base = Path(__file__).resolve().parent
    tpl = base / TEMPLATE_NAME
    out = output_path or (base / OUTPUT_NAME)

    if not tpl.is_file():
        raise FileNotFoundError(f"找不到模板：{tpl}")

    src = Presentation(str(tpl))
    dst = Presentation()
    dst.slide_width = src.slide_width
    dst.slide_height = src.slide_height
    _strip_default_slide(dst)

    if len(SOURCE_SLIDE_ORDER) != len(REPLACEMENTS):
        raise ValueError("SOURCE_SLIDE_ORDER 与 REPLACEMENTS 长度不一致")

    for si, rep in zip(SOURCE_SLIDE_ORDER, REPLACEMENTS):
        slide = duplicate_slide_into(src, si, dst)
        apply_replacements(slide, rep)

    dst.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    p = path.resolve()
    try:
        print(f"已生成：{p}")
    except UnicodeEncodeError:
        print(str(p).encode("utf-8", errors="replace"))
